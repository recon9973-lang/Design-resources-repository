# -*- coding: utf-8 -*-
"""McKinsey 스타일 디자인 시스템 — seulee26/mckinsey-pptx theme.py 토큰 기반.
기존 콘텐츠(build_student/build_instructor)를 그대로 렌더링하도록 동일한 builder API 유지.
액션 타이틀 · 언더룰 · 트래커 · 소스라인/카피라이트 · 네이비 헤더 표 · 신호등 · 가로막대 미감."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SW, SH = Inches(13.333), Inches(7.5)
ML = MR = 0.45
MXin = 0.45
CWin = 13.333 - ML - MR
MX = Inches(ML)
CW = Inches(CWin)

LAT = "Arial"
EA = "Noto Sans CJK KR"
MONO = "Consolas"

# ---- McKinsey palette (theme.py) ----
DEEP = "0A1F3D"; DARK = "0F2A4A"; BRIGHT = "2E9BD6"; MID = "1F6FA8"; LIGHT = "4FB2E5"; ROYAL = "2A2AE5"
INK = "1A1A1A"; BODY = "2B2B2B"; RULE = "999999"; LGRAY = "E8E8E8"; SGRAY = "F2F2F2"; GRID = "D0D0D0"
FOOT = "888888"; PLACE = "BFBFBF"; GREEN = "4CAF50"; AMBER = "E0A44A"; RED = "E04E5E"; WHITE = "FFFFFF"
CODE_BG = "0A1F3D"

# category tags mapped into the McKinsey palette
TAG_COLORS = {
    "웹검색": "1F6FA8",
    "사이트": "2E9BD6",
    "유튜브": "E04E5E",
    "뉴스": "E0A44A",
    "라이브 시연": "0F2A4A",
    "미션": "2A2AE5",
    "실습": "4FB2E5",
    "툴": "1F6FA8",
    "백업": "888888",
}

# old brand-accent hexes -> McKinsey equivalents (so reused content stays on-palette)
REMAP = {
    "B45309": MID, "2563EB": MID, "F59E0B": MID, "4338CA": MID, "C2410C": MID,
    "D97706": MID, "FBBF24": LIGHT, "A5B4FC": LIGHT, "FCD34D": LIGHT, "9A3412": MID,
    "DC2626": RED, "059669": GREEN, "0891B2": LIGHT, "7C3AED": ROYAL, "0F766E": MID,
    "10B981": GREEN, "0EA5E9": BRIGHT, "16A34A": GREEN, "BE185D": ROYAL, "475569": FOOT,
    "14213D": DARK, "292524": DARK,
}


def C(h):
    return RGBColor.from_string(h)


def _rm(color):
    return REMAP.get(color.upper() if isinstance(color, str) else color, color)


def set_font(run, size=12, bold=False, color=BODY, name=LAT, italic=False, spacing=None):
    color = _rm(color)
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = name; f.color.rgb = C(color)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
    ea.set("typeface", EA)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False, radius=None, line_w=0.75):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line); shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded and radius is not None:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp


def add_oval(slide, x, y, d, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    s.shadow.inherit = False; s.fill.solid(); s.fill.fore_color.rgb = C(fill); s.line.fill.background()
    return s


def hline(slide, x1, y, x2, color, wpt=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    ln.line.color.rgb = C(color); ln.line.width = Pt(wpt)
    return ln


def add_textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("sb") is not None: para.space_before = Pt(p["sb"])
        if p.get("sa") is not None: para.space_after = Pt(p["sa"])
        para.line_spacing = p.get("ls", 1.14)
        for txt, opt in p.get("runs", [(p.get("text", ""), {})]):
            r = para.add_run(); r.text = txt
            set_font(r, size=opt.get("size", p.get("size", 12)), bold=opt.get("bold", p.get("bold", False)),
                     color=opt.get("color", p.get("color", BODY)), name=opt.get("name", p.get("name", LAT)),
                     italic=opt.get("italic", p.get("italic", False)), spacing=opt.get("spacing", p.get("spacing")))
    return tb


def new_deck():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_chip(slide, x, y, text, fill, w=None, size=10, color="FFFFFF", h=None):
    h = h or Inches(0.30)
    w = w or Inches(0.24 + 0.125 * max(len(text), 3))
    add_rect(slide, x, y, w, h, fill, rounded=True, radius=0.5)
    add_textbox(slide, x, y + Inches(0.012), w, h,
                [{"text": text, "size": size, "bold": True, "color": color, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
    return w


# ---------------- McKinsey chrome ----------------

def _tracker(slide, label, badge=None):
    w, h, top = 1.85, 0.30, 0.18
    left = 13.333 - MR - w
    add_rect(slide, Inches(left), Inches(top), Inches(w), Inches(h), None, line=PLACE, line_w=0.5)
    add_textbox(slide, Inches(left + 0.08), Inches(top), Inches(w - 0.16), Inches(h),
                [{"runs": [(label, {"size": 10, "color": FOOT})], "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
    if badge:
        bw = 0.9
        add_chip(slide, Inches(left - bw - 0.12), Inches(top), badge, DARK, w=Inches(bw), size=10)


def _footer(slide, theme):
    hline(slide, ML, 7.05, 13.333 - MR, RULE, 0.5)
    src = theme.get("source")
    if src:
        add_textbox(slide, Inches(ML), Inches(7.1), Inches(7.2), Inches(0.22),
                    [{"runs": [("출처  ", {"size": 9, "bold": True, "color": INK}), (src, {"size": 9, "color": INK})]}])
    page = theme.get("_page")
    txt = theme.get("copyright", "") + (f"      {page}" if page else "")
    add_textbox(slide, Inches(13.333 - MR - 6.0), Inches(7.15), Inches(6.0), Inches(0.2),
                [{"runs": [(txt, {"size": 9, "color": FOOT})], "align": PP_ALIGN.RIGHT}])


def s_head(slide, th, title, sub=None, badge=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    kicker = th.get("kicker")
    if kicker:
        add_textbox(slide, Inches(ML), Inches(0.34), Inches(9.6), Inches(0.24),
                    [{"runs": [(kicker, {"size": 10.5, "bold": True, "color": MID, "spacing": 0.8})]}])
    add_textbox(slide, Inches(ML), Inches(0.6), Inches(11.2), Inches(0.64),
                [{"runs": [(title, {"size": 19, "bold": True, "color": INK})], "ls": 1.04}])
    hline(slide, ML, 1.32, 13.333 - MR, RULE, 0.75)
    _tracker(slide, th.get("tracker", ""), badge if badge is not None else th.get("badge"))
    top = 1.5
    if sub:
        add_textbox(slide, Inches(ML), Inches(1.42), Inches(11.4), Inches(0.34),
                    [{"runs": [(sub, {"size": 12, "color": FOOT, "italic": True})], "ls": 1.12}])
        top = 1.86
    _footer(slide, th)
    return top


def s_note(slide, th, note, y=6.42):
    add_rect(slide, MX, Inches(y), CW, Inches(0.44), SGRAY, rounded=False)
    add_rect(slide, MX, Inches(y), Inches(0.06), Inches(0.44), MID)
    add_textbox(slide, MX + Inches(0.2), Inches(y), CW - Inches(0.36), Inches(0.44),
                [{"runs": [("TIP   ", {"size": 10, "bold": True, "color": MID, "spacing": 0.5}),
                           (note, {"size": 10.5, "color": INK})]}], anchor=MSO_ANCHOR.MIDDLE)


# ---------------- builders ----------------

def _bars_motif(slide, x, y):
    vals = [1.4, 2.4, 1.9, 3.0, 2.2]; cols = [MID, BRIGHT, MID, LIGHT, MID]
    for i, (v, cc) in enumerate(zip(vals, cols)):
        add_rect(slide, Inches(x), Inches(y + i * 0.62), Inches(v), Inches(0.34), cc)
    add_rect(slide, Inches(x - 0.02), Inches(y - 0.1), Inches(0.02), Inches(3.2), MID)


def s_title(prs, th, kicker, title, subtitle=None, meta=None, badge=None):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, DEEP)
    _bars_motif(slide, 9.7, 1.7)
    add_rect(slide, Inches(ML + 0.02), Inches(2.15), Inches(0.62), Inches(0.06), BRIGHT)
    add_textbox(slide, Inches(ML + 0.02), Inches(2.36), Inches(9), Inches(0.4),
                [{"runs": [(kicker.upper() if kicker.isascii() else kicker, {"size": 13.5, "bold": True, "color": LIGHT, "spacing": 1.2})]}])
    add_textbox(slide, Inches(ML), Inches(2.95), Inches(8.9), Inches(1.7),
                [{"runs": [(title, {"size": 40, "bold": True, "color": WHITE})], "ls": 1.08}])
    if subtitle:
        add_textbox(slide, Inches(ML + 0.02), Inches(5.15), Inches(8.7), Inches(1.1),
                    [{"runs": [(subtitle, {"size": 15, "color": "C7D3E4"})], "ls": 1.3}])
    hline(slide, ML, 6.72, 13.333 - MR, MID, 0.75)
    if meta:
        add_textbox(slide, Inches(ML), Inches(6.82), Inches(8.5), Inches(0.3),
                    [{"runs": [(meta, {"size": 10, "color": "8DA0BC"})]}])
    b = badge or th.get("badge")
    label = th.get("copyright", "Copyright of AI Marketing Study")
    add_textbox(slide, Inches(13.333 - MR - 4.2), Inches(6.82), Inches(4.2), Inches(0.3),
                [{"runs": [((("[" + b + "]  ") if b else "") + "AI Marketing Study", {"size": 9.5, "color": "6E82A0"})], "align": PP_ALIGN.RIGHT}])
    return slide


def s_big(prs, th, text, sub=None, kicker=None):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, DEEP)
    if kicker:
        add_rect(slide, Inches(6.28), Inches(2.2), Inches(0.62), Inches(0.06), BRIGHT)
        add_textbox(slide, Inches(1.2), Inches(2.4), Inches(10.9), Inches(0.5),
                    [{"runs": [(kicker, {"size": 13.5, "bold": True, "color": LIGHT, "spacing": 1.0})], "align": PP_ALIGN.CENTER}])
    add_textbox(slide, Inches(1.0), Inches(2.98), Inches(11.3), Inches(2.2),
                [{"runs": [(text, {"size": 30, "bold": True, "color": WHITE})], "align": PP_ALIGN.CENTER, "ls": 1.24}])
    if sub:
        add_textbox(slide, Inches(1.7), Inches(5.2), Inches(9.9), Inches(1.4),
                    [{"runs": [(sub, {"size": 15, "color": "C7D3E4"})], "align": PP_ALIGN.CENTER, "ls": 1.34}])
    page = th.get("_page")
    if page:
        add_textbox(slide, Inches(13.333 - MR - 2), Inches(7.05), Inches(2), Inches(0.3),
                    [{"runs": [(str(page), {"size": 9, "color": "6E82A0"})], "align": PP_ALIGN.RIGHT}])
    return slide


def s_bullets(prs, th, title, items, sub=None, note=None, badge=None, size=15):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    paras = []
    for it in items:
        lvl, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        if lvl == 0:
            paras.append({"runs": [("▪  ", {"size": size - 5, "color": MID}),
                                   (txt, {"size": opt.get("size", size), "bold": opt.get("bold", False),
                                          "color": opt.get("color", INK if opt.get("bold") else BODY)})],
                          "sb": 6, "sa": 5, "ls": 1.2})
        elif lvl == 1:
            paras.append({"runs": [("      –  ", {"size": size - 3, "color": PLACE}),
                                   (txt, {"size": opt.get("size", size - 2.5), "bold": opt.get("bold", False),
                                          "color": opt.get("color", BODY)})],
                          "sb": 1, "sa": 2, "ls": 1.16})
        else:
            paras.append({"text": txt, "size": opt.get("size", size + 1), "bold": True, "color": INK, "sb": 12, "sa": 4})
    add_textbox(slide, Inches(ML + 0.05), Inches(top + 0.1), Inches(11.9), Inches(6.3 - top), paras)
    if note:
        s_note(slide, th, note)
    return slide


def _cell(cell, text, size, bold, color, align=PP_ALIGN.LEFT, fill=None):
    if fill:
        cell.fill.solid(); cell.fill.fore_color.rgb = C(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.09)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame; tf.word_wrap = True
    first = True
    for ln in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.line_spacing = 1.08
        r = p.add_run(); r.text = ln
        set_font(r, size=size, bold=bold, color=color)


def s_table(prs, th, title, headers, rows, sub=None, widths=None, note=None,
            badge=None, fsize=12, hsize=12, first_bold=True, tall=False):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    n_r, n_c = len(rows) + 1, len(headers)
    avail_h = (6.34 if note else 6.86) - top
    row_min = 0.40
    gt = slide.shapes.add_table(n_r, n_c, MX, Inches(top + 0.1), CW,
                                Inches(min(avail_h, max(row_min * n_r, 0.44 * n_r))))
    table = gt.table; table.first_row = False; table.horz_banding = False
    if widths:
        total = sum(widths)
        for i, w in enumerate(widths):
            table.columns[i].width = Emu(int(CW * w / total))
    for j, htxt in enumerate(headers):
        _cell(table.cell(0, j), htxt, hsize, True, WHITE, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.LEFT, DARK)
    for i, row in enumerate(rows):
        fill = WHITE if i % 2 == 0 else SGRAY
        for j, val in enumerate(row):
            bold = first_bold and j == 0
            _cell(table.cell(i + 1, j), val, fsize, bold, INK if bold else BODY, PP_ALIGN.LEFT, fill)
    table.rows[0].height = Pt(26)
    for i in range(1, n_r):
        table.rows[i].height = Inches(row_min)
    if note:
        s_note(slide, th, note)
    return slide


def s_cards(prs, th, title, cards, sub=None, cols=2, note=None, badge=None,
            body_size=11.5, head_size=14):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    n = len(cards); rows = (n + cols - 1) // cols
    gap = 0.22; bottom = 6.34 if note else 6.95
    avail_h = bottom - top - 0.08
    ch = (avail_h - gap * (rows - 1)) / rows
    cw = (CWin - gap * (cols - 1)) / cols
    for idx, card in enumerate(cards):
        r, c = divmod(idx, cols)
        x = Inches(ML + c * (cw + gap)); y = Inches(top + 0.1 + r * (ch + gap))
        col = _rm(card.get("tag_color", TAG_COLORS.get(card.get("tag"), MID)))
        add_rect(slide, x, y, Inches(cw), Inches(ch), WHITE, line=GRID, line_w=1.0)
        add_rect(slide, x, y, Inches(cw), Inches(0.05), col)
        cy = y + Inches(0.18); cx = x + Inches(0.2)
        tag = card.get("tag")
        if tag:
            tw = add_chip(slide, cx, cy, tag, col, size=9.5)
            add_textbox(slide, cx + tw + Inches(0.12), cy - Inches(0.02), Inches(cw) - tw - Inches(0.5), Inches(0.4),
                        [{"runs": [(card.get("head", ""), {"size": head_size, "bold": True, "color": INK})]}])
            body_y = cy + Inches(0.42)
        else:
            add_textbox(slide, cx, cy, Inches(cw - 0.38), Inches(0.4),
                        [{"runs": [(card.get("head", ""), {"size": head_size, "bold": True, "color": INK})]}])
            body_y = cy + Inches(0.4)
        paras = []
        for ln in card.get("lines", []):
            if isinstance(ln, tuple):
                label, val = ln
                paras.append({"runs": [(label + "  ", {"size": body_size, "bold": True, "color": MID}),
                                       (val, {"size": body_size, "color": BODY})], "sa": 3, "ls": 1.16})
            else:
                paras.append({"runs": [("· ", {"size": body_size, "color": PLACE}),
                                       (ln, {"size": body_size, "color": BODY})], "sa": 3, "ls": 1.16})
        add_textbox(slide, cx, body_y, Inches(cw - 0.38), y + Inches(ch) - body_y - Inches(0.08), paras)
    if note:
        s_note(slide, th, note)
    return slide


def s_demos(prs, th, title, demos, sub=None, note=None, badge=None):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    n = len(demos); gap = 0.18; bottom = 6.34 if note else 6.95
    avail = bottom - top - 0.08; ch = (avail - gap * (n - 1)) / n; q_h = 0.42
    for i, d in enumerate(demos):
        y = top + 0.1 + i * (ch + gap); yI = Inches(y)
        add_rect(slide, MX, yI, CW, Inches(ch), WHITE, line=GRID, line_w=1.0)
        tag = d.get("tag", ""); tcol = _rm(TAG_COLORS.get(tag, MID))
        add_rect(slide, MX, yI, Inches(0.08), Inches(ch), tcol)
        rx = MX + Inches(0.24)
        add_textbox(slide, rx, Inches(y + 0.16), Inches(1.55), Inches(0.5),
                    [{"runs": [(d.get("when", ""), {"size": 12, "bold": True, "color": INK})], "ls": 1.1}])
        add_chip(slide, rx, Inches(y + ch - 0.46), tag, tcol, size=9.5)
        add_rect(slide, MX + Inches(1.82), Inches(y + 0.16), Inches(0.012), Inches(ch - 0.32), LGRAY)
        bx = MX + Inches(2.02); bw = CW - Inches(2.24)
        add_textbox(slide, bx, Inches(y + 0.13), bw, Inches(0.5),
                    [{"runs": [(d.get("what", ""), {"size": 13, "bold": True, "color": INK})], "ls": 1.12}])
        qy = y + 0.13 + 0.4
        add_rect(slide, bx, Inches(qy), bw, Inches(q_h), SGRAY)
        add_textbox(slide, bx + Inches(0.12), Inches(qy + 0.02), bw - Inches(0.24), Inches(q_h - 0.04),
                    [{"runs": [("검색/링크   ", {"size": 9, "bold": True, "color": MID}),
                               (d.get("q", ""), {"size": 10, "color": BODY, "name": MONO})], "ls": 1.08}], anchor=MSO_ANCHOR.MIDDLE)
        my = qy + q_h + 0.08
        add_textbox(slide, bx, Inches(my), bw, Inches(y + ch - my - 0.06),
                    [{"runs": [("멘트   ", {"size": 9.5, "bold": True, "color": tcol}),
                               (d.get("say", ""), {"size": 11, "color": BODY})], "ls": 1.14}])
    if note:
        s_note(slide, th, note)
    return slide


def s_prompt(prs, th, title, code, sub=None, note=None, badge=None, size=11.5):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    bottom = 6.34 if note else 6.9
    box_y = top + 0.1
    add_rect(slide, MX, Inches(box_y), CW, Inches(bottom - box_y), CODE_BG)
    add_rect(slide, MX, Inches(box_y), CW, Inches(0.34), DARK)
    for i, dot in enumerate(["FF5F56", "FFBD2E", "27C93F"]):
        add_oval(slide, MX + Inches(0.2 + i * 0.22), Inches(box_y + 0.12), Inches(0.1), dot)
    add_textbox(slide, MX + Inches(0.95), Inches(box_y), CW - Inches(1.1), Inches(0.34),
                [{"runs": [("PROMPT", {"size": 9.5, "bold": True, "color": "7FA8CC", "spacing": 1.2})]}], anchor=MSO_ANCHOR.MIDDLE)
    paras = []
    for line in code.strip("\n").split("\n"):
        color = "E5ECF5"
        if line.strip().endswith(":") and len(line) < 30: color = "7DD3FC"
        if line.startswith("#"): color = "7FA8CC"
        paras.append({"text": line if line else " ", "size": size, "color": color, "name": MONO, "sa": 2, "ls": 1.16})
    add_textbox(slide, MX + Inches(0.3), Inches(box_y + 0.5), CW - Inches(0.6), Inches(bottom - box_y - 0.62), paras)
    if note:
        s_note(slide, th, note)
    return slide


def s_twocol(prs, th, title, left, right, sub=None, note=None, badge=None, body_size=13):
    slide = blank(prs)
    top = s_head(slide, th, title, sub, badge)
    bottom = 6.34 if note else 6.9
    h = bottom - top - 0.1; w = (CWin - 0.26) / 2
    for i, col in enumerate([left, right]):
        x = Inches(ML + i * (w + 0.26)); y = Inches(top + 0.1)
        hd = _rm(col.get("color", MID if i == 0 else DARK))
        add_rect(slide, x, y, Inches(w), Inches(h), WHITE, line=GRID, line_w=1.0)
        add_rect(slide, x, y, Inches(w), Inches(0.5), hd)
        add_textbox(slide, x + Inches(0.2), y, Inches(w - 0.4), Inches(0.5),
                    [{"runs": [(col["head"], {"size": 14.5, "bold": True, "color": WHITE})]}], anchor=MSO_ANCHOR.MIDDLE)
        paras = []
        for it in col["items"]:
            lvl, txt = (it[0], it[1]) if isinstance(it, tuple) else (0, it)
            if lvl == 0:
                paras.append({"runs": [("▪  ", {"size": body_size - 5, "color": hd}),
                                       (txt, {"size": body_size, "color": INK})], "sb": 5, "sa": 3, "ls": 1.18})
            else:
                paras.append({"runs": [("      – ", {"size": body_size - 3, "color": PLACE}),
                                       (txt, {"size": body_size - 1.5, "color": BODY})], "sa": 2, "ls": 1.14})
        add_textbox(slide, x + Inches(0.22), y + Inches(0.66), Inches(w - 0.44), Inches(h - 0.8), paras)
    if note:
        s_note(slide, th, note)
    return slide


BUILDERS = {"title": s_title, "big": s_big, "bullets": s_bullets, "table": s_table,
            "cards": s_cards, "prompt": s_prompt, "twocol": s_twocol, "demos": s_demos}


def build_deck(path, theme, slides):
    prs = new_deck()
    for i, (name, kwargs) in enumerate(slides):
        theme["_page"] = i + 1
        BUILDERS[name](prs, theme, **kwargs)
    prs.save(path)
    return len(slides)
