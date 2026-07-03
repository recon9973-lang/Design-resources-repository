# -*- coding: utf-8 -*-
"""PPTX 디자인 시스템 v2 — "발견·인용 / Signal & Citation".

개선 초점(디자인 리뷰 반영):
- 가독성 우선: 표 본문 최소 12pt, 행 높이/여백 확대, 명확한 zebra.
- 밀집한 강사용 D1~D5 표 → 풀폭 '시연 카드'로 재구성(별도 builder: demos).
- hue-bias 중립색, 절제된 강조색, 학생/강사 브랜드 구분.
- 타이틀 '시그널 바' 모티프, 카드/프롬프트/팁/2단 컴포넌트 정돈.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SW, SH = Inches(13.333), Inches(7.5)
MX = Inches(0.62)
CW = Inches(12.093)

FONT = "맑은 고딕"
MONO = "Consolas"

# ---- shared neutrals (hue-biased, not flat grey) ----
INK = "111927"
BODY = "374151"
MUTE = "6B7280"
FAINT = "9AA4B2"
PAPER = "F6F7F9"
CARD = "FFFFFF"
LINE = "E3E8EF"
CODE_BG = "0E1526"

# demo/category tags — semantic colors (encode real categories)
TAG_COLORS = {
    "웹검색": "2563EB",
    "사이트": "0F766E",
    "유튜브": "DC2626",
    "뉴스": "B45309",
    "라이브 시연": "7C3AED",
    "미션": "BE185D",
    "실습": "0891B2",
    "툴": "4338CA",
    "백업": "475569",
}


def C(h):
    return RGBColor.from_string(h)


def set_font(run, size=14, bold=False, color=BODY, name=FONT, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = name
    f.color.rgb = C(color)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False, radius=None, line_w=0.75):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    if line:
        shp.line.color.rgb = C(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded and radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("sb") is not None:
            para.space_before = Pt(p["sb"])
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        para.line_spacing = p.get("ls", 1.14)
        runs = p.get("runs")
        if runs is None:
            runs = [(p.get("text", ""), {})]
        for txt, opt in runs:
            r = para.add_run()
            r.text = txt
            set_font(
                r,
                size=opt.get("size", p.get("size", 14)),
                bold=opt.get("bold", p.get("bold", False)),
                color=opt.get("color", p.get("color", BODY)),
                name=opt.get("name", p.get("name", FONT)),
                italic=opt.get("italic", p.get("italic", False)),
                spacing=opt.get("spacing", p.get("spacing")),
            )
    return tb


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_chip(slide, x, y, text, fill, w=None, size=10.5, color="FFFFFF", h=None):
    h = h or Inches(0.30)
    w = w or Inches(0.24 + 0.132 * max(len(text), 3))
    add_rect(slide, x, y, w, h, fill, rounded=True, radius=0.5)
    add_textbox(
        slide, x, y + Inches(0.012), w, h,
        [{"text": text, "size": size, "bold": True, "color": color, "align": PP_ALIGN.CENTER,
          "spacing": 0.2}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return w


def add_footer(slide, th):
    add_rect(slide, MX, Inches(7.16), Inches(0.20), Inches(0.028), th["accent"])
    add_textbox(
        slide, MX + Inches(0.30), Inches(7.05), CW, Inches(0.3),
        [{"runs": [(th.get("footer", ""), {"size": 9, "color": FAINT})]}],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def s_head(slide, th, title, sub=None, badge=None):
    add_rect(slide, 0, 0, SW, Inches(0.14), th["accent"])
    add_rect(slide, MX, Inches(0.50), Inches(0.09), Inches(0.50), th["accent"], rounded=True, radius=0.5)
    add_textbox(
        slide, MX + Inches(0.26), Inches(0.42), Inches(10.0), Inches(0.66),
        [{"text": title, "size": 22, "bold": True, "color": INK}],
    )
    badge = badge if badge is not None else th.get("badge")
    if badge:
        wch = Inches(0.42 + 0.145 * len(badge))
        add_chip(slide, SW - wch - Inches(0.62), Inches(0.52), badge, th["accent"], w=wch, size=11)
    top = 1.20
    if sub:
        add_textbox(
            slide, MX + Inches(0.27), Inches(1.02), Inches(11.4), Inches(0.4),
            [{"text": sub, "size": 12.5, "color": MUTE, "ls": 1.15}],
        )
        top = 1.52
    add_footer(slide, th)
    return top


def s_note(slide, th, note, y=6.60):
    add_rect(slide, MX, Inches(y), CW, Inches(0.46), th["soft"], rounded=True, radius=0.22)
    add_rect(slide, MX, Inches(y), Inches(0.075), Inches(0.46), th["accent"])
    add_textbox(
        slide, MX + Inches(0.24), Inches(y), CW - Inches(0.42), Inches(0.46),
        [{"runs": [("TIP   ", {"size": 10.5, "bold": True, "color": th["accent"], "spacing": 0.5}),
                   (note, {"size": 11, "color": BODY})]}],
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ---------------- slide builders ----------------

def _signal_bars(slide, x, y, color, heights):
    bw = Inches(0.052)
    gap = Inches(0.042)
    for i, hh in enumerate(heights):
        add_rect(slide, x + i * (bw + gap), y + Inches(0.34 - hh), bw, Inches(hh), color, rounded=True, radius=0.4)


def s_title(prs, th, kicker, title, subtitle=None, meta=None, badge=None):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, th["ground"])
    add_rect(slide, 0, SH - Inches(0.16), SW, Inches(0.16), th["accent"])
    # signal-bars motif (measurement vernacular)
    _signal_bars(slide, MX + Inches(0.02), Inches(1.98), th["accent_lt"],
                 [0.16, 0.27, 0.20, 0.34, 0.24, 0.30, 0.14])
    add_textbox(
        slide, MX + Inches(0.02), Inches(2.44), Inches(11.5), Inches(0.5),
        [{"text": kicker, "size": 14.5, "bold": True, "color": th["accent_lt"], "spacing": 0.6}],
    )
    add_textbox(
        slide, MX + Inches(0.02), Inches(2.90), Inches(12.0), Inches(1.9),
        [{"text": title, "size": 39, "bold": True, "color": "FFFFFF", "ls": 1.06}],
    )
    if subtitle:
        add_rect(slide, MX + Inches(0.05), Inches(4.66), Inches(0.5), Inches(0.035), th["accent"])
        add_textbox(
            slide, MX + Inches(0.02), Inches(4.86), Inches(11.5), Inches(1.2),
            [{"text": subtitle, "size": 16, "color": "D5DAE6", "ls": 1.32}],
        )
    if meta:
        add_textbox(
            slide, MX + Inches(0.02), Inches(6.66), Inches(11.6), Inches(0.5),
            [{"text": meta, "size": 11, "color": th["meta"]}],
        )
    b = badge or th.get("badge")
    if b:
        wch = Inches(0.5 + 0.16 * len(b))
        add_rect(slide, SW - wch - Inches(0.62), Inches(0.6), wch, Inches(0.44), None, line=th["accent_lt"], rounded=True, radius=0.5, line_w=1.25)
        add_textbox(
            slide, SW - wch - Inches(0.62), Inches(0.6), wch, Inches(0.44),
            [{"text": b, "size": 12.5, "bold": True, "color": th["accent_lt"], "align": PP_ALIGN.CENTER, "spacing": 0.5}],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    return slide


def s_big(prs, th, text, sub=None, kicker=None):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, th["ground"])
    add_rect(slide, 0, SH - Inches(0.16), SW, Inches(0.16), th["accent"])
    if kicker:
        _signal_bars(slide, Inches(6.28), Inches(1.86), th["accent_lt"], [0.14, 0.24, 0.18, 0.28, 0.16])
        add_textbox(
            slide, Inches(1.2), Inches(2.34), Inches(10.9), Inches(0.5),
            [{"text": kicker, "size": 14.5, "bold": True, "color": th["accent_lt"], "align": PP_ALIGN.CENTER, "spacing": 0.6}],
        )
    add_textbox(
        slide, Inches(1.0), Inches(2.92), Inches(11.3), Inches(2.2),
        [{"text": text, "size": 30, "bold": True, "color": "FFFFFF", "align": PP_ALIGN.CENTER, "ls": 1.24}],
    )
    if sub:
        add_textbox(
            slide, Inches(1.7), Inches(5.18), Inches(9.9), Inches(1.4),
            [{"text": sub, "size": 15, "color": "CBD2E0", "align": PP_ALIGN.CENTER, "ls": 1.36}],
        )
    return slide


def s_bullets(prs, th, title, items, sub=None, note=None, badge=None, size=15.5):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    paras = []
    for it in items:
        lvl, txt = it[0], it[1]
        opt = it[2] if len(it) > 2 else {}
        if lvl == 0:
            paras.append({
                "runs": [("■  ", {"size": size - 6, "color": th["accent"]}),
                         (txt, {"size": opt.get("size", size), "bold": opt.get("bold", False),
                                "color": opt.get("color", INK if opt.get("bold") else BODY)})],
                "sb": 7, "sa": 5, "ls": 1.2,
            })
        elif lvl == 1:
            paras.append({
                "runs": [("      –  ", {"size": size - 3, "color": FAINT}),
                         (txt, {"size": opt.get("size", size - 2.5), "bold": opt.get("bold", False),
                                "color": opt.get("color", BODY)})],
                "sb": 1, "sa": 2, "ls": 1.16,
            })
        else:
            paras.append({
                "text": txt, "size": opt.get("size", size + 1), "bold": True,
                "color": opt.get("color", INK), "sb": 12, "sa": 4,
            })
    add_textbox(slide, MX + Inches(0.16), Inches(top + 0.14), Inches(11.6), Inches(6.35 - top), paras)
    if note:
        s_note(slide, th, note)
    return slide


def _fill_cell(cell, text, size, bold, color, align=PP_ALIGN.LEFT, fill=None):
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = C(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.11)
    cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.045)
    cell.margin_bottom = Inches(0.045)
    tf = cell.text_frame
    tf.word_wrap = True
    lines = str(text).split("\n")
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = ln
        set_font(r, size=size, bold=bold, color=color)


def _no_table_border(table):
    # remove default banding style artifacts by setting explicit fills only
    pass


def s_table(prs, th, title, headers, rows, sub=None, widths=None, note=None,
            badge=None, fsize=12, hsize=12, first_bold=True, tall=False):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    n_r, n_c = len(rows) + 1, len(headers)
    avail_h = (6.42 if note else 6.86) - top
    row_min = 0.40
    gt = slide.shapes.add_table(n_r, n_c, MX, Inches(top + 0.12), CW,
                                Inches(min(avail_h, max(row_min * n_r, 0.44 * n_r))))
    table = gt.table
    table.first_row = False
    table.horz_banding = False
    if widths:
        total = sum(widths)
        for i, w in enumerate(widths):
            table.columns[i].width = Emu(int(CW * w / total))
    for j, htxt in enumerate(headers):
        _fill_cell(table.cell(0, j), htxt, hsize, True, "FFFFFF",
                   PP_ALIGN.CENTER, th["accent"])
    for i, row in enumerate(rows):
        fill = CARD if i % 2 == 0 else th["soft"]
        for j, val in enumerate(row):
            bold = first_bold and j == 0
            _fill_cell(table.cell(i + 1, j), val, fsize, bold,
                       INK if bold else BODY, PP_ALIGN.LEFT, fill)
    table.rows[0].height = Pt(26)
    for i in range(1, n_r):
        table.rows[i].height = Inches(row_min)
    if note:
        s_note(slide, th, note)
    return slide


def s_cards(prs, th, title, cards, sub=None, cols=2, note=None, badge=None,
            body_size=11.5, head_size=14.5):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    n = len(cards)
    rows = (n + cols - 1) // cols
    gap = 0.24
    bottom = 6.5 if note else 7.0
    avail_h = bottom - top - 0.1
    ch = (avail_h - gap * (rows - 1)) / rows
    cw = (12.093 - gap * (cols - 1)) / cols
    for idx, card in enumerate(cards):
        r, c = divmod(idx, cols)
        x = Inches(0.62 + c * (cw + gap))
        y = Inches(top + 0.12 + r * (ch + gap))
        add_rect(slide, x, y, Inches(cw), Inches(ch), CARD, line=LINE, rounded=True, radius=0.05, line_w=1.0)
        add_rect(slide, x, y, Inches(0.07), Inches(ch), card.get("tag_color", TAG_COLORS.get(card.get("tag"), th["accent"])), rounded=False)
        cy = y + Inches(0.16)
        cx = x + Inches(0.24)
        tag = card.get("tag")
        if tag:
            tw = add_chip(slide, cx, cy, tag, card.get("tag_color", TAG_COLORS.get(tag, th["accent"])), size=10)
            add_textbox(
                slide, cx + tw + Inches(0.14), cy - Inches(0.02), Inches(cw) - tw - Inches(0.55), Inches(0.4),
                [{"text": card.get("head", ""), "size": head_size, "bold": True, "color": INK}],
            )
            body_y = cy + Inches(0.44)
        else:
            add_textbox(
                slide, cx, cy, Inches(cw - 0.42), Inches(0.4),
                [{"text": card.get("head", ""), "size": head_size, "bold": True, "color": INK}],
            )
            body_y = cy + Inches(0.42)
        paras = []
        for ln in card.get("lines", []):
            if isinstance(ln, tuple):
                label, val = ln
                paras.append({"runs": [(label + "  ", {"size": body_size, "bold": True, "color": th["accent"]}),
                                       (val, {"size": body_size, "color": BODY})],
                              "sa": 3.5, "ls": 1.18})
            else:
                paras.append({"runs": [("· ", {"size": body_size, "color": FAINT}),
                                       (ln, {"size": body_size, "color": BODY})],
                              "sa": 3.5, "ls": 1.18})
        add_textbox(slide, cx, body_y, Inches(cw - 0.42), y + Inches(ch) - body_y - Inches(0.1), paras)
    if note:
        s_note(slide, th, note)
    return slide


def s_demos(prs, th, title, demos, sub=None, note=None, badge=None):
    """풀폭 시연 카드: 각 카드 = 시점배지 + 유형태그 + 무엇/검색어(mono)/멘트."""
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    n = len(demos)
    gap = 0.20
    bottom = 6.48 if note else 7.0
    avail = bottom - top - 0.12
    ch = (avail - gap * (n - 1)) / n
    q_h = 0.44  # mono 검색/링크 영역(최대 2줄)
    for i, d in enumerate(demos):
        y = top + 0.14 + i * (ch + gap)
        yI = Inches(y)
        add_rect(slide, MX, yI, CW, Inches(ch), CARD, line=LINE, rounded=True, radius=0.045, line_w=1.0)
        tag = d.get("tag", "")
        tcol = TAG_COLORS.get(tag, th["accent"])
        add_rect(slide, MX, yI, Inches(0.09), Inches(ch), tcol)
        # left rail: 시점 + 유형
        rx = MX + Inches(0.26)
        add_textbox(slide, rx, Inches(y + 0.16), Inches(1.5), Inches(0.5),
                    [{"text": d.get("when", ""), "size": 12, "bold": True, "color": INK, "ls": 1.1}])
        add_chip(slide, rx, Inches(y + ch - 0.46), tag, tcol, size=10)
        # divider
        add_rect(slide, MX + Inches(1.80), Inches(y + 0.16), Inches(0.014), Inches(ch - 0.32), LINE)
        # body
        bx = MX + Inches(2.02)
        bw = CW - Inches(2.24)
        add_textbox(slide, bx, Inches(y + 0.13), bw, Inches(0.5),
                    [{"runs": [(d.get("what", ""), {"size": 13, "bold": True, "color": INK})], "ls": 1.14}])
        qy = y + 0.13 + 0.40
        add_rect(slide, bx, Inches(qy), bw, Inches(q_h), th["soft"], rounded=True, radius=0.14)
        add_textbox(slide, bx + Inches(0.14), Inches(qy + 0.02), bw - Inches(0.28), Inches(q_h - 0.04),
                    [{"runs": [("검색/링크   ", {"size": 9, "bold": True, "color": th["accent"], "spacing": 0.3}),
                               (d.get("q", ""), {"size": 10, "color": BODY, "name": MONO})], "ls": 1.1}],
                    anchor=MSO_ANCHOR.MIDDLE)
        my = qy + q_h + 0.09
        add_textbox(slide, bx, Inches(my), bw, Inches(y + ch - my - 0.08),
                    [{"runs": [("멘트   ", {"size": 10, "bold": True, "color": tcol, "spacing": 0.3}),
                               (d.get("say", ""), {"size": 11, "color": BODY})], "ls": 1.16}])
    if note:
        s_note(slide, th, note)
    return slide


def s_prompt(prs, th, title, code, sub=None, note=None, badge=None, size=11.5):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    bottom = 6.46 if note else 6.92
    box_y = top + 0.12
    add_rect(slide, MX, Inches(box_y), CW, Inches(bottom - box_y), CODE_BG, rounded=True, radius=0.025)
    # header strip
    add_rect(slide, MX, Inches(box_y), CW, Inches(0.36), "1B2440", rounded=True, radius=0.025)
    add_rect(slide, MX, Inches(box_y + 0.18), CW, Inches(0.18), "1B2440")
    for i, dot in enumerate(["FF5F56", "FFBD2E", "27C93F"]):
        add_rect(slide, MX + Inches(0.24 + i * 0.22), Inches(box_y + 0.13), Inches(0.10), Inches(0.10), dot, rounded=True, radius=0.5)
    add_textbox(slide, MX + Inches(1.0), Inches(box_y), CW - Inches(1.2), Inches(0.36),
                [{"text": "PROMPT", "size": 9.5, "bold": True, "color": "7C8AA5", "spacing": 1.2}],
                anchor=MSO_ANCHOR.MIDDLE)
    paras = []
    for line in code.strip("\n").split("\n"):
        color = "E5E9F0"
        if line.strip().endswith(":") and len(line) < 30:
            color = "7DD3FC"
        if line.startswith("#"):
            color = "8391A8"
        paras.append({"text": line if line else " ", "size": size, "color": color,
                      "name": MONO, "sa": 2, "ls": 1.16})
    add_textbox(slide, MX + Inches(0.34), Inches(box_y + 0.52), CW - Inches(0.68),
                Inches(bottom - box_y - 0.66), paras)
    if note:
        s_note(slide, th, note)
    return slide


def s_twocol(prs, th, title, left, right, sub=None, note=None, badge=None, body_size=13):
    slide = blank(prs)
    add_rect(slide, 0, 0, SW, SH, PAPER)
    top = s_head(slide, th, title, sub, badge)
    bottom = 6.46 if note else 6.92
    h = bottom - top - 0.12
    w = (12.093 - 0.28) / 2
    for i, col in enumerate([left, right]):
        x = Inches(0.62 + i * (w + 0.28))
        y = Inches(top + 0.12)
        hd_color = col.get("color", th["accent"] if i == 0 else col.get("color2", th["accent2"]))
        add_rect(slide, x, y, Inches(w), Inches(h), CARD, line=LINE, rounded=True, radius=0.035, line_w=1.0)
        add_rect(slide, x, y, Inches(w), Inches(0.54), hd_color, rounded=True, radius=0.045)
        add_rect(slide, x, y + Inches(0.27), Inches(w), Inches(0.27), hd_color)
        add_textbox(slide, x + Inches(0.22), y, Inches(w - 0.44), Inches(0.54),
                    [{"text": col["head"], "size": 15, "bold": True, "color": "FFFFFF"}],
                    anchor=MSO_ANCHOR.MIDDLE)
        paras = []
        for it in col["items"]:
            lvl, txt = (it[0], it[1]) if isinstance(it, tuple) else (0, it)
            if lvl == 0:
                paras.append({"runs": [("■  ", {"size": body_size - 6, "color": hd_color}),
                                       (txt, {"size": body_size, "color": INK})],
                              "sb": 6, "sa": 3, "ls": 1.2})
            else:
                paras.append({"runs": [("      – ", {"size": body_size - 3, "color": FAINT}),
                                       (txt, {"size": body_size - 1.5, "color": BODY})],
                              "sa": 2, "ls": 1.15})
        add_textbox(slide, x + Inches(0.24), y + Inches(0.72), Inches(w - 0.48), Inches(h - 0.9), paras)
    if note:
        s_note(slide, th, note)
    return slide


BUILDERS = {
    "title": s_title,
    "big": s_big,
    "bullets": s_bullets,
    "table": s_table,
    "cards": s_cards,
    "prompt": s_prompt,
    "twocol": s_twocol,
    "demos": s_demos,
}


def build_deck(path, theme, slides):
    prs = new_deck()
    for name, kwargs in slides:
        BUILDERS[name](prs, theme, **kwargs)
    prs.save(path)
    return len(slides)


# ---- brand themes ----
THEME_STUDENT = {
    "ground": "161E3A",     # deep indigo ground (title/section/statement)
    "accent": "4338CA",     # indigo primary
    "accent_lt": "A5B4FC",  # light indigo (on dark)
    "accent2": "F59E0B",    # amber highlight (used sparingly)
    "soft": "EEF0FB",       # indigo-tinted zebra/soft panel
    "meta": "8A94B2",
    "badge": None,
}

THEME_INSTRUCTOR = {
    "ground": "2A211C",     # espresso ground
    "accent": "C2410C",     # burnt clay primary
    "accent_lt": "FDBA74",  # light clay (on dark)
    "accent2": "9A3412",
    "soft": "FBEDE4",       # clay-tinted zebra/soft panel
    "meta": "B7A99E",
    "badge": "강사용",
}
